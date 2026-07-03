"""Create an editable PowerPoint brief from one website audit."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "outputs"


def _rgb(value: str, fallback: str) -> RGBColor:
    raw = (value or fallback).lstrip("#")
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    return RGBColor.from_string(raw.upper())


def _set_run(run, *, font: str, size: float, color: RGBColor, bold: bool = False, italic: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_title_slide(prs: Presentation, audit: dict, display_font: str, body_font: str, white: RGBColor, accent: RGBColor, brand: RGBColor) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = brand

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10.8), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "WEBSITE REVIEW"
    _set_run(run, font=body_font, size=10, color=accent, bold=True)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11), Inches(2.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = audit["brand"]["name"]
    _set_run(run, font=display_font, size=28, color=white, bold=False)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = audit["client_review"]["headline"]
    _set_run(run, font=body_font, size=16, color=white)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"{audit['url']} | {audit['generated_at'][:10]}"
    _set_run(run, font=body_font, size=9, color=white)


def _add_text_slide(
    prs: Presentation,
    *,
    kicker: str,
    title: str,
    lines: list[str],
    display_font: str,
    body_font: str,
    ink: RGBColor,
    muted: RGBColor,
    accent: RGBColor,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("FBF9F6")

    kicker_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(4.5), Inches(0.4))
    p = kicker_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = kicker.upper()
    _set_run(run, font=body_font, size=9, color=accent, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.0), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, font=display_font, size=24, color=ink)

    body_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(11.0), Inches(5.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = line
        _set_run(run, font=body_font, size=14, color=muted)


def _add_two_column_slide(
    prs: Presentation,
    *,
    kicker: str,
    title: str,
    left_title: str,
    left_lines: list[str],
    right_title: str,
    right_lines: list[str],
    display_font: str,
    body_font: str,
    ink: RGBColor,
    muted: RGBColor,
    accent: RGBColor,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("FBF9F6")

    kicker_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(4.5), Inches(0.4))
    p = kicker_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = kicker.upper()
    _set_run(run, font=body_font, size=9, color=accent, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, font=display_font, size=22, color=ink)

    for x, heading, lines in (
        (0.8, left_title, left_lines),
        (6.4, right_title, right_lines),
    ):
        heading_box = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(5.0), Inches(0.4))
        p = heading_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = heading
        _set_run(run, font=body_font, size=11, color=accent, bold=True)

        body_box = slide.shapes.add_textbox(Inches(x), Inches(2.35), Inches(5.1), Inches(4.6))
        tf = body_box.text_frame
        tf.word_wrap = True
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = line
            _set_run(run, font=body_font, size=13, color=muted)


def render_editable_pptx(audit: dict, report_id: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUTPUT_DIR / f"{report_id}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    display_font = audit["brand"].get("display_font") or "Georgia"
    body_font = audit["brand"].get("body_font") or "Arial"
    ink = _rgb("#1a1a1a", "#1a1a1a")
    muted = _rgb("#625f5b", "#625f5b")
    white = _rgb("#ffffff", "#ffffff")
    brand = _rgb(audit["brand"].get("primary_hex"), "#1a1a1a")
    accent = _rgb(audit["brand"].get("accent_hex"), "#b89a6a")

    _add_title_slide(prs, audit, display_font, body_font, white, accent, brand)
    _add_text_slide(
        prs,
        kicker="The short answer",
        title=audit["client_review"]["headline"],
        lines=[
            f"Current conversion readiness: {audit['conversion']['score']}/100.",
            f"Primary root layer: {audit['conversion']['root_layer']}.",
            "The site already has a clear point of view. The next gain comes from making the choice, proof, and next step easier.",
        ],
        display_font=display_font,
        body_font=body_font,
        ink=ink,
        muted=muted,
        accent=accent,
    )
    _add_text_slide(
        prs,
        kicker="What already works",
        title="Keep the parts that help people understand the business.",
        lines=audit["client_review"]["strengths"] + [
            f"{item['title']}: {item['text']}" for item in audit["client_review"]["visitor_view"]
        ],
        display_font=display_font,
        body_font=body_font,
        ink=ink,
        muted=muted,
        accent=accent,
    )
    _add_two_column_slide(
        prs,
        kicker="Copy review",
        title="The message recognizes the problem. It needs more proof and fewer abstract words.",
        left_title="What works",
        left_lines=audit["client_review"]["copy"]["works"],
        right_title="What weakens it",
        right_lines=audit["client_review"]["copy"]["risks"],
        display_font=display_font,
        body_font=body_font,
        ink=ink,
        muted=muted,
        accent=accent,
    )
    _add_text_slide(
        prs,
        kicker="What would make this stronger",
        title="Every score comes with the next move.",
        lines=[
            f"{item['plain_name']} — {item['score']}/{item['max']} — {item['to_10']}"
            for item in audit["client_review"]["conversion_path"]
        ],
        display_font=display_font,
        body_font=body_font,
        ink=ink,
        muted=muted,
        accent=accent,
    )

    if audit["mode"] in {"visibility", "full"}:
        _add_text_slide(
            prs,
            kicker="Search and AI visibility",
            title="Being easy to quote is not the same as being easy to understand.",
            lines=[
                audit["client_review"]["visibility_explanation"],
                *[
                    f"{item['name']} — {item['score']}/{item['max']} — {item['to_10']}"
                    for item in audit["client_review"]["visibility_path"]
                ],
            ],
            display_font=display_font,
            body_font=body_font,
            ink=ink,
            muted=muted,
            accent=accent,
        )

    _add_text_slide(
        prs,
        kicker="First week",
        title="Start with the changes that make the buying decision easier.",
        lines=[f"{index}. {item}" for index, item in enumerate(audit["seven_day_plan"][:5], start=1)],
        display_font=display_font,
        body_font=body_font,
        ink=ink,
        muted=muted,
        accent=accent,
    )

    prs.core_properties.title = f"{audit['brand']['name']} Website Review"
    prs.core_properties.subject = "Editable website audit deck"
    prs.core_properties.author = "Edgewise"
    prs.save(path)
    return path
